#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
V.Together 과제 결과 보고서 PPT 생성 스크립트
실행: python3 scripts/generate-result-ppt.py
생성 파일: docs/V.Together_과제결과보고서.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_PATH = Path(__file__).resolve().parent.parent / "docs" / "V.Together_과제결과보고서.pptx"

# ───── 색상 팔레트 ─────
GREEN_DARK   = RGBColor(0x16, 0x65, 0x34)
GREEN_MID    = RGBColor(0x15, 0x80, 0x3D)
GREEN_LIGHT  = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_ACCENT = RGBColor(0x22, 0xC5, 0x5E)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x1F, 0x29, 0x37)
GRAY_700     = RGBColor(0x37, 0x41, 0x51)
GRAY_500     = RGBColor(0x6B, 0x72, 0x80)
GRAY_300     = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_100     = RGBColor(0xF3, 0xF4, 0xF6)
BLUE_600     = RGBColor(0x25, 0x63, 0xEB)
AMBER_600    = RGBColor(0xD9, 0x77, 0x06)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.8)
CONTENT_W = Inches(11.7)


def _fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill_shape(shape, color)
    shape.line.fill.background()
    return shape


def _set_para(paragraph, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.alignment = align


def _add_text(slide, left, top, width, height, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_para(p, text, size, bold, color, align)
    p.line_spacing = line_spacing
    return box


def _add_multiline(slide, left, top, width, height, lines, size=14, color=BLACK, line_spacing=1.15, bullet=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (txt, is_bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.bold = is_bold
        p.font.color.rgb = color
        p.line_spacing = line_spacing
        p.space_before = Pt(4) if i > 0 else Pt(0)
    return box


def _slide_base(prs, title_text, subtitle_text=None):
    """공통 슬라이드: 상단 초록 바 + 제목"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), GREEN_DARK)
    _add_text(slide, MARGIN_L, Inches(0.22), CONTENT_W, Inches(0.7),
              title_text, size=26, bold=True, color=WHITE)
    if subtitle_text:
        _add_text(slide, MARGIN_L, Inches(0.72), CONTENT_W, Inches(0.35),
                  subtitle_text, size=13, color=RGBColor(0xBB, 0xF7, 0xD0))
    _add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.06), GREEN_ACCENT)
    return slide


def _add_card(slide, left, top, width, height, title, body_lines, title_color=GREEN_DARK):
    """카드 모양 (둥근 테두리는 pptx에서 제한적이므로 사각 + 좌측 선)"""
    bg = _add_rect(slide, left, top, width, height, GRAY_100)
    bar = _add_rect(slide, left, top, Inches(0.06), height, title_color)

    _add_text(slide, left + Inches(0.25), top + Inches(0.12), width - Inches(0.4), Inches(0.35),
              title, size=13, bold=True, color=title_color)
    body_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.48), width - Inches(0.4), height - Inches(0.55))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY_700
        p.line_spacing = 1.25
        p.space_before = Pt(2) if i > 0 else Pt(0)


def _add_table(slide, left, top, width, rows_data, col_widths=None):
    """표 추가. rows_data = [[cell, ...], ...] 첫 행 = 헤더"""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.45 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE if r == 0 else GRAY_700
            p.font.bold = (r == 0)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_DARK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else GRAY_100
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


# ══════════════════════════════════════════
# 슬라이드 생성
# ══════════════════════════════════════════
def build(prs):
    # ─── 1. 표지 ───
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(cover, 0, 0, SLIDE_W, SLIDE_H, GREEN_DARK)
    _add_rect(cover, 0, Inches(3.0), SLIDE_W, Inches(0.08), GREEN_ACCENT)

    _add_text(cover, MARGIN_L, Inches(1.8), CONTENT_W, Inches(1),
              "V.Together", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(cover, MARGIN_L, Inches(2.7), CONTENT_W, Inches(0.6),
              "ESG 임직원 참여 플랫폼", size=20, color=RGBColor(0xBB, 0xF7, 0xD0), align=PP_ALIGN.CENTER)
    _add_text(cover, MARGIN_L, Inches(3.6), CONTENT_W, Inches(0.8),
              "과제 결과 보고서", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(cover, MARGIN_L, Inches(5.8), CONTENT_W, Inches(0.4),
              "2026년 3월  |  VNTG", size=16, color=RGBColor(0x86, 0xEF, 0xAC), align=PP_ALIGN.CENTER)

    # ─── 2. 목차 ───
    s = _slide_base(prs, "목차", "Table of Contents")
    items = [
        ("01", "프로젝트 개요"),
        ("02", "시스템 아키텍처 및 기술 스택"),
        ("03", "데이터 모델"),
        ("04", "구현 기능 요약 (Phase 1~3)"),
        ("05", "사용자 화면 — 메인·기부·마이"),
        ("06", "이벤트 & 챌린지 시스템"),
        ("07", "관리자 기능"),
        ("08", "인증 방식 상세 (객관식 CHOICE 포함)"),
        ("09", "보상 체계 및 매칭 정책"),
        ("10", "V.Honors & MAU 지표"),
        ("11", "개발 산출물 및 마이그레이션"),
        ("12", "결론 및 향후 계획"),
    ]
    for i, (num, title) in enumerate(items):
        y = Inches(1.55) + Inches(i * 0.44)
        _add_text(s, MARGIN_L, y, Inches(0.6), Inches(0.35),
                  num, size=16, bold=True, color=GREEN_ACCENT)
        _add_text(s, Inches(1.5), y, Inches(8), Inches(0.35),
                  title, size=16, color=BLACK)

    # ─── 3. 프로젝트 개요 ───
    s = _slide_base(prs, "01  프로젝트 개요", "Project Overview")
    _add_multiline(s, MARGIN_L, Inches(1.5), CONTENT_W, Inches(1.5), [
        ("배경", True),
        ("ESG 경영 실천 및 조직문화(Culture) 활성화를 위한 통합 임직원 참여 플랫폼 구축", False),
        ("VNTG 전직원(약 230명) 대상, PC/Mobile 반응형 웹 서비스", False),
        ("사내 G-Suite(@vntg.co.kr) Google OAuth 2.0 인증", False),
    ], size=14, color=GRAY_700)

    _add_card(s, MARGIN_L, Inches(3.3), Inches(3.6), Inches(1.4),
              "핵심 목표 1", ["전사 연간 기부 목표", "2,000만 포인트 달성"])
    _add_card(s, Inches(4.8), Inches(3.3), Inches(3.6), Inches(1.4),
              "핵심 목표 2", ["전직원 80% 이상이 월 1회 이상", "접속하는 MAU 확보"])
    _add_card(s, Inches(8.8), Inches(3.3), Inches(3.6), Inches(1.4),
              "인증 방식", ["Google OAuth 2.0", "사내 G-Suite 계정 전용"])

    _add_card(s, MARGIN_L, Inches(5.1), Inches(3.6), Inches(1.5),
              "Gamification", ["ESG Level(3단계), 뱃지, 랭킹", "동기 부여 시스템"], title_color=BLUE_600)
    _add_card(s, Inches(4.8), Inches(5.1), Inches(3.6), Inches(1.5),
              "Transparency", ["기부금 모금 현황 실시간 공개", "달성률 Progress Bar"], title_color=BLUE_600)
    _add_card(s, Inches(8.8), Inches(5.1), Inches(3.6), Inches(1.5),
              "Togetherness", ["동료 칭찬, 팀 대항전", "사내 연대감 강화"], title_color=BLUE_600)

    # ─── 4. 시스템 아키텍처 ───
    s = _slide_base(prs, "02  시스템 아키텍처 및 기술 스택", "Technical Architecture")
    _add_table(s, MARGIN_L, Inches(1.5), Inches(11.5), [
        ["구분", "기술", "용도"],
        ["프레임워크", "Next.js 16 (App Router)", "프론트/백 통합, React 19, SSR/SSG"],
        ["BaaS", "Supabase", "Auth, PostgreSQL DB, Storage, RLS"],
        ["API 전략", "Server Actions ('use server')", "데이터 변경(INSERT/UPDATE/DELETE)"],
        ["UI 라이브러리", "Shadcn UI + Tailwind CSS", "디자인 시스템, 반응형(Mobile First)"],
        ["상태 관리", "Zustand + React Hooks", "전역 상태 / 로컬 상태"],
        ["리치 텍스트", "TipTap Editor", "이벤트 상세 설명 편집"],
        ["엑셀", "xlsx (SheetJS)", "제출 목록 .xlsx 다운로드"],
        ["배포", "Google Cloud Run", "서버리스 컨테이너 호스팅"],
    ], col_widths=[Inches(1.8), Inches(3.5), Inches(6.2)])

    # ─── 5. 데이터 모델 ───
    s = _slide_base(prs, "03  데이터 모델", "ERD — 핵심 테이블")
    _add_table(s, MARGIN_L, Inches(1.5), Inches(11.5), [
        ["테이블", "주요 컬럼", "설명"],
        ["users", "user_id, dept_name, current_points, total_donated_amount, level, is_admin, last_active_at", "사용자 정보, ESG 레벨, 관리자 여부"],
        ["donation_targets", "target_amount, current_amount, status", "기부처별 목표·현재 모금액"],
        ["donations", "user_id, target_id, amount, created_at", "개별 기부 트랜잭션"],
        ["events", "type(SEASONAL/ALWAYS), category, status, reward_policy, frequency_limit", "이벤트/챌린지 마스터"],
        ["event_rounds", "event_id, round_number, start_date, end_date, submission_deadline", "기간제 이벤트의 구간 정보"],
        ["event_submissions", "event_id, user_id, round_id, status, reward_received, is_anonymous", "참여자 인증 제출·심사 상태"],
        ["event_verification_methods", "event_id, method_type, input_style, label, instruction, unit, options", "인증 방식 설정 (사진/텍스트/숫자/동료선택/객관식)"],
        ["event_rewards", "event_id, reward_type, amount", "보상 유형 (V.Credit/굿즈/커피쿠폰)"],
        ["point_transactions", "user_id, amount, type, related_type, user_email, user_name", "포인트 원장 (적립/차감 내역)"],
    ], col_widths=[Inches(2.6), Inches(5.5), Inches(3.4)])

    # ─── 6. 구현 기능 요약 ───
    s = _slide_base(prs, "04  구현 기능 요약 (Phase 1 ~ 3)", "Development Progress")

    _add_card(s, MARGIN_L, Inches(1.5), Inches(3.6), Inches(2.6),
              "Phase 1: 기반 구축 & 기부 (MVP)", [
                  "[완료] 사용자 DB · OAuth 연동",
                  "[완료] 메인 대시보드(레벨, 포인트)",
                  "[완료] 기부 기능(목표/마감 로직)",
                  "[완료] 관리자 기본 페이지",
                  "  - 사용자 목록, 포인트 지급",
                  "  - 메인 문구 편집, 관리자 토글",
              ], title_color=GREEN_MID)

    _add_card(s, Inches(4.8), Inches(1.5), Inches(3.6), Inches(2.6),
              "Phase 2: 챌린지 & 보상 시스템", [
                  "[완료] 이벤트 목록/등록/수정",
                  "[완료] 사용자 참여 모달(동적 폼)",
                  "[완료] 칭찬 챌린지 쌍방 지급",
                  "[완료] 인증 심사 센터(대량 처리)",
                  "[완료] 보상 선택(CHOICE)",
                  "[완료] 쿠폰/굿즈 발송 관리",
                  "[완료] 엑셀 다운로드",
              ], title_color=BLUE_600)

    _add_card(s, Inches(8.8), Inches(1.5), Inches(3.6), Inches(2.6),
              "Phase 3: 고도화 & 안정화", [
                  "[완료] V.Honors TOP 10 랭킹",
                  "[완료] 모바일 UX 최적화",
                  "[완료] MAU 지표",
                  "[완료] 관리자 UX 고도화",
                  "[선택] Redis 캐싱",
                  "[선택] Vertex AI 이미지 분류",
              ], title_color=AMBER_600)

    _add_card(s, MARGIN_L, Inches(4.5), Inches(11.7), Inches(2.5),
              "추가 완료 항목 상세", [
                  "• V.Credit 수동 지급 (/admin/point-grant) — 이벤트 외 보정·특별 보상 시 직원에게 P 직접 적립, 사유 기록(ADMIN_GRANT)",
                  "• ALWAYS 이벤트 빈도 제한 (ONCE/DAILY/WEEKLY/MONTHLY) — 제출 전 검사 + 모달 안내",
                  "• 반려 알림 — 카드/상세에 반려 표시, 마이페이지에서 반려 사유 확인",
                  "• 칭찬 챌린지 익명 옵션 — 수신자에게 '익명의 동료'로 표시, 관리자는 실명 확인 가능",
                  "• 세아웍스 인사 연동 — 로그인 시 부서명(dept_name) 자동 동기화",
                  "• Soft Delete — 모든 테이블 deleted_at 컬럼, 실제 삭제 없이 플래그 처리",
              ])

    # ─── 7. 사용자 화면 ───
    s = _slide_base(prs, "05  사용자 화면", "메인 · 기부 · 마이페이지")

    _add_card(s, MARGIN_L, Inches(1.5), Inches(5.6), Inches(2.4),
              "메인 대시보드  /", [
                  "My Status: ESG Level, 보유 포인트, 누적 기부액",
                  "전사 기부 목표 달성률 Progress Bar",
                  "이벤트 & 챌린지 카드 (전체/V.Together/People 탭)",
                  "V.Honors 명예의 전당 TOP 10 (개인/팀)",
              ])
    _add_card(s, Inches(6.8), Inches(1.5), Inches(5.6), Inches(2.4),
              "기부 페이지  /donation", [
                  "기부처 카드 (4~5개): 목표/현재 모금액, Progress",
                  "달성 완료 시: Gold 테두리, 뱃지, 버튼 비활성화",
                  "기부하기 모달: 보유 포인트 내에서 자유 금액 입력",
                  "개인 기부(급여 공제): 외부 Google Form 연결",
              ])
    _add_card(s, MARGIN_L, Inches(4.3), Inches(5.6), Inches(2.4),
              "마이페이지  /my", [
                  "프로필(이름, 부서, ESG Level 뱃지)",
                  "포인트 내역: 적립/사용/기부 이력",
                  "참여 이벤트 내역: 상태(대기/승인/반려) + 반려 사유",
                  "보상 선택 대기: 미수령 보상 표시",
              ])
    _add_card(s, Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.4),
              "ESG Level 시스템", [
                  "Eco Keeper (LV1):  0 ~ 30,000 P   회색 뱃지",
                  "Green Master (LV2): 30,001 ~ 80,000 P  초록 뱃지",
                  "Earth Hero (LV3):  80,001 P ~   보라 뱃지",
                  "기준: 누적 기부액(total_donated_amount)",
              ])

    # ─── 8. 이벤트 & 챌린지 ───
    s = _slide_base(prs, "06  이벤트 & 챌린지 시스템", "운영 방식 · 참여 · 인증 흐름")

    _add_table(s, MARGIN_L, Inches(1.5), Inches(11.5), [
        ["이벤트 타입", "운영 방식", "참여 제한", "예시"],
        ["SEASONAL (기간제)", "구간별 기간 설정, 인증 마감일 별도", "구간당 1회", "월간 텀블러 챌린지"],
        ["ALWAYS (상시)", "기간 없음, 빈도 제한만", "ONCE/DAILY/WEEKLY/MONTHLY", "일회성 설문, 매주 운동"],
        ["INTERACTIVE (칭찬)", "동료 선택 + 메시지, 쌍방 보상", "빈도 제한에 따름", "칭찬 챌린지 (익명 가능)"],
    ], col_widths=[Inches(2.2), Inches(3.5), Inches(3.0), Inches(2.8)])

    _add_multiline(s, MARGIN_L, Inches(3.7), CONTENT_W, Inches(3.5), [
        ("이벤트 상태 관리", True),
        ("• ACTIVE: 메인 화면 노출, 참여 가능  |  PAUSED: 일시 정지, 노출 안 됨  |  ENDED: 종료 (관리자에서 재개 가능)", False),
        ("", False),
        ("참여 흐름 (사용자 → 관리자)", True),
        ("1. 사용자: 이벤트 카드에서 [인증하기] 클릭 → 동적 폼(사진/텍스트/숫자/동료선택/객관식) 작성 후 제출", False),
        ("2. 제출 상태: 항상 PENDING으로 저장 (자동 승인 없음)", False),
        ("3. 관리자: /admin/verifications 에서 내용 확인 후 승인 또는 반려 (일괄 처리 지원)", False),
        ("4. 승인 → 보상 지급 (단일 V.Credit: 즉시 지급 / 복수 보상: 사용자 선택 대기)", False),
        ("", False),
        ("기간제 구간 구조 (월 3구간)", True),
        ("• 1구간: 1~10일 (인증 마감 15일)  |  2구간: 11~20일 (인증 마감 25일)  |  3구간: 21~말일 (인증 마감 익월 5일)", False),
    ], size=13, color=GRAY_700)

    # ─── 9. 관리자 기능 ───
    s = _slide_base(prs, "07  관리자 기능", "/admin — 대시보드, 이벤트, 심사, 포인트, 기부처")

    _add_table(s, MARGIN_L, Inches(1.5), Inches(11.5), [
        ["페이지", "URL", "주요 기능"],
        ["관리자 대시보드", "/admin", "전사 기부, 목표 달성률, 승인 대기, MAU, 이벤트 적립 현황(People/V.Together/매칭금)"],
        ["이벤트 목록", "/admin/events", "전체 이벤트 목록, 상태 토글(종료/재개), 엑셀 다운로드"],
        ["이벤트 등록", "/admin/events/new", "보상(V.Credit/굿즈/커피쿠폰) 복수 선택, 인증 방식 다중 설정, 구간 관리"],
        ["이벤트 상세·수정", "/admin/events/[id]", "소개문구(TipTap), 상태, 보상 금액, 구간 추가/삭제"],
        ["인증 심사 센터", "/admin/verifications", "대량 승인/반려, 사진 썸네일·텍스트 미리보기, 필터·검색"],
        ["V.Credit 수동 지급", "/admin/point-grant", "이벤트 외 보정·특별 보상, 사유 기록(ADMIN_GRANT)"],
        ["기부처 관리", "/admin/donation-targets", "목표 수정, 오프라인 성금 합산"],
        ["쿠폰/굿즈 발송", "/admin/reward-fulfillment", "발송 대상 목록, 발송 완료 체크, 필터(전체/미발송/발송완료)"],
        ["최근 접속", "/admin/recent-users", "마지막 접속 시각 기준 사용자 목록"],
    ], col_widths=[Inches(2.0), Inches(2.8), Inches(6.7)])

    _add_multiline(s, MARGIN_L, Inches(5.8), CONTENT_W, Inches(1.2), [
        ("관리자 UX 고도화 (2026.03.18)", True),
        ("네비 아이콘·승인 대기 배지 | 설정 섹션 접기/펼치기 | 인증 심사 기본 필터·검색 | AdminPageHeader·breadcrumb | 기부처 표시명 통일", False),
    ], size=12, color=GRAY_500)

    # ─── 10. 인증 방식 상세 ───
    s = _slide_base(prs, "08  인증 방식 상세", "사진 · 텍스트 · 숫자 · 동료선택 · 객관식(CHOICE)")

    _add_table(s, MARGIN_L, Inches(1.5), Inches(11.5), [
        ["인증 타입", "input_style", "설명", "UI"],
        ["PHOTO (사진)", "—", "Supabase Storage 업로드, 2장 이상 필수", "파일 업로드 폼"],
        ["TEXT (텍스트)", "SHORT (단답)", "한 줄 입력", "input 필드"],
        ["TEXT (텍스트)", "LONG (장문)", "여러 줄 입력", "textarea 필드"],
        ["TEXT (텍스트)", "CHOICE (객관식)", "관리자 정의 선택지 중 택 1", "라디오/드롭다운"],
        ["VALUE (숫자)", "—", "숫자만 입력, 단위(km, km/h 등) 표시", "number input + 단위"],
        ["PEER_SELECT", "—", "동료 검색/선택 (칭찬 챌린지)", "사용자 검색 폼"],
    ], col_widths=[Inches(2.2), Inches(2.0), Inches(4.5), Inches(2.8)])

    _add_multiline(s, MARGIN_L, Inches(4.6), CONTENT_W, Inches(2.5), [
        ("마이그레이션 031: 객관식(CHOICE) 추가", True),
        ("", False),
        ("• event_verification_methods.input_style CHECK 제약에 'CHOICE' 추가", False),
        ("• event_verification_methods.options JSONB 컬럼 추가 — 선택지 문자열 배열 저장 (예: [\"A\", \"B\", \"C\"])", False),
        ("• 텍스트(TEXT) 항목에서만 단답(SHORT) / 장문(LONG) / 객관식(CHOICE) 선택 가능", False),
        ("• 모든 인증 방식에 label(제목), instruction(안내 문구) 공통 지원", False),
        ("• 항목 여러 개 조합 가능 (예: 사진 + 텍스트, 동료선택 + 텍스트, 텍스트 2개 등)", False),
    ], size=13, color=GRAY_700)

    # ─── 11. 보상 체계 ───
    s = _slide_base(prs, "09  보상 체계 및 매칭 정책", "V.Credit · 굿즈 · 커피쿠폰 · People 매칭")

    _add_card(s, MARGIN_L, Inches(1.5), Inches(3.6), Inches(2.0),
              "V.Credit (포인트)", [
                  "단일 보상 시: 승인 즉시 자동 지급",
                  "복수 보상 시: 사용자 선택 후 지급",
                  "current_points 증가 + 거래 기록",
              ], title_color=GREEN_MID)
    _add_card(s, Inches(4.8), Inches(1.5), Inches(3.6), Inches(2.0),
              "커피 쿠폰", [
                  "amount(금액) 필수 설정",
                  "승인 후 관리자 별도 발송",
                  "/admin/reward-fulfillment 관리",
              ], title_color=AMBER_600)
    _add_card(s, Inches(8.8), Inches(1.5), Inches(3.6), Inches(2.0),
              "굿즈", [
                  "금액 없음 (실물 보상)",
                  "승인 후 관리자 별도 발송",
                  "발송 완료 체크 + 필터",
              ], title_color=BLUE_600)

    _add_multiline(s, MARGIN_L, Inches(3.9), CONTENT_W, Inches(3.2), [
        ("보상 선택 (CHOICE) 흐름", True),
        ("1. 이벤트에 보상 2종 이상 등록 (예: V.Credit + 커피쿠폰)", False),
        ("2. 관리자 승인 → reward_received=false, 사용자 선택 대기", False),
        ("3. 사용자: 메인에서 [보상받기] 클릭 → 모달에서 원하는 보상 선택", False),
        ("4. claimRewardChoice 서버 액션으로 지급 완료", False),
        ("", False),
        ("카테고리별 매칭 정책", True),
        ("• People 이벤트: 사용자 V.Credit 적립액만큼 회사 매칭 (1:1)", False),
        ("• V.Together 이벤트: 매칭 없음", False),
        ("• 전체 모인금액 = V.Together 적립 + People 적립 + People 매칭금액", False),
        ("• 관리자 대시보드에서 People/V.Together별 현황 확인 가능", False),
    ], size=13, color=GRAY_700)

    # ─── 12. V.Honors & MAU ───
    s = _slide_base(prs, "10  V.Honors & MAU 지표", "명예의 전당 · 월간 활성 사용자")

    _add_card(s, MARGIN_L, Inches(1.5), Inches(5.6), Inches(2.8),
              "V.Honors (명예의 전당)", [
                  "메인 페이지에만 TOP 10 표시 (전용 페이지 없음)",
                  "개인 랭킹 / 팀 랭킹 탭",
                  "분기(Q1~Q4)별 리셋 — 해당 분기 기부액 기준",
                  "1~3위 메달 표시, 레벨 뱃지 컬러 적용",
                  "누적 기부액은 별도 유지 (ESG Level용)",
                  "",
                  "모바일: 가로 스크롤, 터치 영역 44px+",
              ])
    _add_card(s, Inches(6.8), Inches(1.5), Inches(5.6), Inches(2.8),
              "MAU (월간 활성 사용자)", [
                  "정의: 최근 30일 내 1회 이상 접속한 고유 사용자 수",
                  "",
                  "수집 방식:",
                  "  users.last_active_at 컬럼",
                  "  getCurrentUser() 호출 시마다 갱신",
                  "",
                  "표시: 관리자 대시보드 'MAU (최근 30일)' 카드",
                  "마이그레이션: 016-users-last-active-at.sql",
              ])

    # ─── 13. 산출물 ───
    s = _slide_base(prs, "11  개발 산출물 및 마이그레이션", "소스 코드 · DB 마이그레이션 · 문서")

    _add_card(s, MARGIN_L, Inches(1.5), Inches(5.6), Inches(2.3),
              "소스 코드 구조", [
                  "app/          페이지 라우팅 (13개 페이지)",
                  "api/actions/   서버 액션 (데이터 변경)",
                  "api/queries/   데이터 조회 함수",
                  "components/    공통 UI 컴포넌트",
                  "hooks/         커스텀 React Hooks",
                  "store/         Zustand 전역 상태",
                  "lib/           유틸리티 (rounds, seah-orgsync 등)",
              ])
    _add_card(s, Inches(6.8), Inches(1.5), Inches(5.6), Inches(2.3),
              "DB 마이그레이션 (29개)", [
                  "docs/migrations/ 폴더에 버전 관리",
                  "002 ~ 031 순차 실행 (Supabase SQL Editor)",
                  "주요: 006(이벤트 테이블), 014(input_style),",
                  "016(MAU), 017(쿠폰 발송), 020(Soft Delete),",
                  "031(객관식 CHOICE + options JSONB)",
              ])
    _add_card(s, MARGIN_L, Inches(4.2), Inches(5.6), Inches(2.3),
              "문서 (docs/)", [
                  "PRD.md — 마스터 기획서 (v2.6)",
                  "logic.md — 비즈니스 로직 정리",
                  "progress.md — 개발 진행 현황",
                  "plan-admin.md — 관리자 페이지 설계서",
                  "plan-events-operations.md — 이벤트 운영 방식",
                  "plan-phase3.md — Phase 3 기술 설계서",
              ])
    _add_card(s, Inches(6.8), Inches(4.2), Inches(5.6), Inches(2.3),
              "페이지 URL 목록 (13개)", [
                  "/            메인 대시보드",
                  "/donation     기부",
                  "/my           마이페이지",
                  "/login        로그인",
                  "/admin        관리자 대시보드 외 8개 서브 페이지",
              ])

    # ─── 14. 결론 ───
    s = _slide_base(prs, "12  결론 및 향후 계획", "Summary & Future Plan")

    _add_card(s, MARGIN_L, Inches(1.5), Inches(11.7), Inches(2.0),
              "주요 성과", [
                  "• Phase 1~3 핵심 기능 전량 구현 완료 — 기부, 이벤트/챌린지, 인증 심사, V.Credit 지급, 보상 선택, 관리자 대시보드",
                  "• 인증 입력 형태에 객관식(CHOICE) 추가하여 이벤트 운영 유연성 향상 (마이그레이션 031)",
                  "• 관리자 UX 고도화: 네비 아이콘, 승인 대기 배지, 검색/필터, breadcrumb 등 운영 편의성 대폭 개선",
                  "• 세아웍스 인사 연동으로 부서 정보 자동 동기화, Soft Delete 적용으로 데이터 안전성 확보",
              ], title_color=GREEN_MID)

    _add_card(s, MARGIN_L, Inches(3.9), Inches(5.6), Inches(2.2),
              "향후 계획 (선택적 고도화)", [
                  "V.Honors Redis 캐싱 — 인프라 도입 시",
                  "Vertex AI 이미지 분류 — 인증 사진 자동 검증",
                  "프로필 이미지 연동 — 세아웍스 API 필드 추가 시",
                  "이용약관/개인정보처리방침 페이지 연동",
              ], title_color=GRAY_500)

    _add_card(s, Inches(6.8), Inches(3.9), Inches(5.6), Inches(2.2),
              "기술 특징", [
                  "Server Components 우선 설계",
                  "Server Actions로 보안·간결한 API",
                  "모바일 퍼스트 반응형 UI",
                  "29개 마이그레이션으로 스키마 버전 관리",
                  "RLS(Row Level Security) 적용",
              ], title_color=GRAY_500)

    # ─── 15. 끝 (Thank you) ───
    end = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(end, 0, 0, SLIDE_W, SLIDE_H, GREEN_DARK)
    _add_text(end, MARGIN_L, Inches(2.5), CONTENT_W, Inches(1),
              "감사합니다", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_rect(end, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.06), GREEN_ACCENT)
    _add_text(end, MARGIN_L, Inches(4.2), CONTENT_W, Inches(0.6),
              "V.Together — ESG 임직원 참여 플랫폼", size=18, color=RGBColor(0xBB, 0xF7, 0xD0), align=PP_ALIGN.CENTER)
    _add_text(end, MARGIN_L, Inches(5.0), CONTENT_W, Inches(0.4),
              "VNTG  |  2026", size=14, color=RGBColor(0x86, 0xEF, 0xAC), align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)
    build(prs)
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"생성 완료: {OUTPUT_PATH}")
    print(f"총 슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    main()
